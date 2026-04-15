"""HTML to PPTX conversion helpers using python-pptx."""

from __future__ import annotations

import io
import logging
import re
from html import unescape

from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


def _strip_tags(html: str) -> str:
    """Remove HTML tags and return plain text."""
    text = re.sub(r"<br\s*/?>", "\n", html, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", "", text)
    return unescape(text).strip()


def _split_into_sections(html: str) -> list[tuple[str, str]]:
    """Split HTML into (heading, body) sections for slides."""
    # Split on h1/h2 headings
    parts = re.split(r"(<h[12][^>]*>.*?</h[12]>)", html, flags=re.IGNORECASE | re.DOTALL)

    sections: list[tuple[str, str]] = []
    current_heading = ""
    current_body = ""

    for part in parts:
        part = part.strip()
        if not part:
            continue
        if re.match(r"<h[12]", part, re.IGNORECASE):
            if current_heading or current_body:
                sections.append((current_heading, current_body.strip()))
            current_heading = _strip_tags(part)
            current_body = ""
        else:
            current_body += " " + _strip_tags(part)

    if current_heading or current_body:
        sections.append((current_heading, current_body.strip()))

    # If no sections found, put everything in one slide
    if not sections:
        text = _strip_tags(html)
        if text:
            sections.append(("Document", text))

    return sections


def html_to_pptx_bytes(html: str, title: str = "Document") -> bytes:
    """Convert HTML string to PPTX bytes."""
    from app.utils.sanitization import sanitize_html_content

    safe_html = sanitize_html_content(html) or ""
    sections = _split_into_sections(safe_html)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = f"{len(sections)} sections"

    # Content slides
    content_layout = prs.slide_layouts[1]
    for heading, body in sections:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = heading or "Untitled Section"
        # Truncate body to avoid huge slides
        display_body = body[:2000] + ("..." if len(body) > 2000 else "")
        if slide.placeholders[1]:
            tf = slide.placeholders[1].text_frame
            tf.text = display_body
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
