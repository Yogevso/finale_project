"""PDF-to-PPTX converter.

Renders each PDF page to an image and places it on its own PowerPoint slide.
The generated PPTX is used as a preserved working attachment for PDF uploads.
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Emu


_PAGE_RENDER_DPI = 144


@dataclass
class PdfToPptxConversionResult:
    """Result of converting PDF bytes into PPTX bytes."""

    pptx_bytes: bytes = b""
    page_count: int = 0
    warnings: list[str] = field(default_factory=list)
    error: str | None = None


def convert_pdf_to_pptx(pdf_bytes: bytes) -> PdfToPptxConversionResult:
    """Convert raw PDF bytes into a slide deck with one slide per page."""

    result = PdfToPptxConversionResult()
    try:
        pdf = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception as exc:  # policy: FAIL_FAST — invalid PDF input returns a stable conversion error
        result.error = f"Failed to open PDF: {exc}"
        return result

    result.page_count = len(pdf)
    if result.page_count == 0:
        result.error = "PDF has no pages"
        pdf.close()
        return result

    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    for page_idx in range(result.page_count):
        page = pdf[page_idx]
        try:
            pix = page.get_pixmap(dpi=_PAGE_RENDER_DPI, alpha=False)
            image_bytes = pix.tobytes("png")
        except Exception as exc:  # policy: FAIL_FAST — page render failure aborts conversion cleanly
            result.error = f"Failed to render PDF page {page_idx + 1}: {exc}"
            pdf.close()
            return result

        slide = presentation.slides.add_slide(blank_layout)
        try:
            _add_page_image_to_slide(
                slide,
                image_bytes=image_bytes,
                slide_width=slide_width,
                slide_height=slide_height,
                warnings=result.warnings,
                page_idx=page_idx,
            )
        except Exception as exc:  # policy: FAIL_FAST — generated slide deck must be structurally valid
            result.error = f"Failed to build PPTX slide for PDF page {page_idx + 1}: {exc}"
            pdf.close()
            return result

    pdf.close()

    output = io.BytesIO()
    presentation.save(output)
    result.pptx_bytes = output.getvalue()
    return result


def _add_page_image_to_slide(
    slide,
    *,
    image_bytes: bytes,
    slide_width: int,
    slide_height: int,
    warnings: list[str],
    page_idx: int,
) -> None:
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            width_px, height_px = image.size
    except Exception as exc:  # policy: FAIL_FAST — rendered page metadata must be readable
        raise ValueError(f"Failed to inspect rendered PDF page image: {exc}") from exc

    if width_px <= 0 or height_px <= 0:
        raise ValueError("Rendered PDF page has invalid dimensions")

    scale = min(slide_width / width_px, slide_height / height_px)
    picture_width = Emu(int(width_px * scale))
    picture_height = Emu(int(height_px * scale))
    left = Emu(int((slide_width - picture_width) / 2))
    top = Emu(int((slide_height - picture_height) / 2))

    try:
        slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            left,
            top,
            width=picture_width,
            height=picture_height,
        )
    except Exception as exc:  # policy: FAIL_FAST — slide image placement must succeed or conversion fails
        warnings.append(f"Page {page_idx + 1}: failed to place rendered image on slide: {exc}")
        raise
